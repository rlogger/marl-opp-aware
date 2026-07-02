"""Build the research talk (.pptx) with figures + full speaker notes.

Reproducible python-pptx builder. scripts/build_deck.mjs produced the previous
13-slide artifact but depends on @oai/artifact-tool, which is not on npm; this
builder keeps that deck's arc and wording and extends it with the slides the
6/12+ meeting asks require (occupancy heatmaps, scatter interpretation, the
averaging caveat, the BC failure demo, plot-generation procedures).

Output: marl_opp_aware_results.pptx (repo root).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PLOTS = os.path.join(ROOT, "plots")
INK = RGBColor(0x17, 0x17, 0x17)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)
ACCENT = RGBColor(0x0D, 0x6E, 0x7A)
WARNING = RGBColor(0xB4, 0x42, 0x33)
EMW, EMH = 13.333, 7.5


def _fit(path, maxw, maxh):
    w, h = Image.open(path).size
    ar = w / h
    iw, ih = maxw, maxw / ar
    if ih > maxh:
        ih, iw = maxh, maxh * ar
    return iw, ih


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(EMW)
        self.prs.slide_height = Inches(EMH)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, notes):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        s.notes_slide.notes_text_frame.text = notes
        if self.n > 1:
            tb = s.shapes.add_textbox(Inches(EMW - 0.9), Inches(EMH - 0.5),
                                      Inches(0.6), Inches(0.3))
            p = tb.text_frame.paragraphs[0]
            p.text = str(self.n)
            p.font.size = Pt(11)
            p.font.color.rgb = MUTED
        return s

    def text(self, s, txt, left, top, w, h, size=22, bold=False,
             color=MUTED, center=False):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(txt.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
            if center:
                from pptx.enum.text import PP_ALIGN
                p.alignment = PP_ALIGN.CENTER
        return tb

    def title(self, s, main, kicker=None):
        if kicker:
            self.text(s, kicker.upper(), 0.55, 0.28, 4, 0.35, size=12,
                      bold=True, color=ACCENT)
        self.text(s, main, 0.55, 0.55, EMW - 1.1, 0.9, size=27, bold=True,
                  color=INK)

    def bullets(self, s, items, left, top, w, h, size=16):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(size)
            p.font.color.rgb = MUTED
            p.space_after = Pt(8)

    def img(self, s, name, left, top, maxw, maxh):
        path = os.path.join(PLOTS, name)
        iw, ih = _fit(path, maxw, maxh)
        s.shapes.add_picture(path, Inches(left + (maxw - iw) / 2),
                             Inches(top + (maxh - ih) / 2),
                             Inches(iw), Inches(ih))

    def number(self, s, num, caption, left, top, color=ACCENT):
        self.text(s, num, left, top, 3.2, 0.7, size=40, bold=True, color=color)
        self.text(s, caption, left, top + 0.72, 3.2, 0.9, size=13, color=MUTED)

    def plot_slide(self, title, kicker, plot, bullets, callout, notes):
        s = self.slide(notes)
        self.title(s, title, kicker)
        self.img(s, plot, 0.5, 1.5, 7.4, 5.4)
        self.bullets(s, bullets, 8.2, 1.8, 4.6, 3.4, size=16)
        self.text(s, callout, 8.2, 5.3, 4.6, 1.5, size=14, bold=True, color=INK)


def main():
    d = Deck()

    # 1 · title (matches the prior artifact: centered title, otherwise blank)
    s = d.slide("Open slowly. The whole talk is one controlled question: can a "
                "team measure an opponent's hidden strategy, keep uncertainty "
                "in the loop, and use that belief during planning?")
    d.text(s, "Adaptive Opponent\nModeling for MARL", 1.6, 2.5, EMW - 3.2, 1.9,
           size=48, bold=True, color=INK, center=True)

    # 2 · claim
    s = d.slide("Use this slide to state the exact claim. The evidence is "
                "captures per episode, not action accuracy. 4.31 is the main "
                "positive result. 1.42 explains why hard labels can hurt "
                "control. 4.08 shows the belief survives removing intent labels.")
    d.title(s, "The paper claim is narrow enough to defend", "Claim")
    d.text(s, "A predator team performs better when it plans against a "
              "calibrated belief over the prey's hidden intent.",
           0.6, 1.6, 9.5, 1.1, size=24, bold=True, color=INK)
    d.number(s, "4.31", "captures/episode with inferred-belief planning", 0.7, 3.4)
    d.number(s, "1.42", "captures/episode with a wrong hard intent", 4.9, 3.4, WARNING)
    d.number(s, "4.08", "captures/episode with the label-free JEPA belief", 9.1, 3.4, INK)

    # 3 · setup
    s = d.slide("Lived detail: three blue predators chase one red prey in the "
                "JaxMARL simple-tag arena; the red prey privately draws one of "
                "four corner intents at reset. The prey knows the intent; the "
                "predators infer it from motion. The map is fixed, so any lift "
                "must come from reading behavior, which isolates opponent "
                "modeling from map perception.")
    d.title(s, "The task hides strategy inside the prey", "Setup")
    rows = [("Episode draw", "The prey receives one hidden corner intent at reset."),
            ("Observation", "Predators see motion, positions, obstacles, and rewards."),
            ("Belief", "The first k prey steps produce a posterior over four intents."),
            ("Control", "The planner samples that posterior during lookahead.")]
    y = 1.7
    for label, body in rows:
        d.text(s, label, 0.9, y, 2.6, 0.5, size=18, bold=True, color=INK)
        d.text(s, body, 3.9, y, 8.6, 0.6, size=18)
        y += 0.85
    d.text(s, "The map stays fixed. Any lift must come from reading behavior "
              "and using the belief.", 0.9, y + 0.3, 11.4, 0.6, size=18,
           bold=True, color=INK)

    # 4 · demo gif (task in motion)
    s = d.slide("Play this in slideshow mode; the GIF animates. Four episodes, "
                "one per corner. Left panel: intent-blind predators hedge. "
                "Right: belief-driven predators infer the corner and intercept; "
                "the X is the JEPA-predicted destination and the bar strip is "
                "the belief sharpening. Narrate one episode: early scramble, "
                "belief flat; prey commits, belief spikes, pack converges.")
    d.title(s, "The task in motion: hedge vs infer-and-intercept", "Demo")
    d.img(s, "demo_jepa.gif", 1.2, 1.45, 11.0, 5.6)

    # 5 · part 1 (belief measurable)
    d.plot_slide(
        "The belief signal becomes measurable", "Part 1", "part2_intent_eval.png",
        ["Accuracy rises from 0.37 at 3 observed steps to 0.97 at 25.",
         "Posterior entropy falls from 1.35 to 0.03 nats.",
         "The belief sharpens before the episode ends."],
        "This validates the Part 1 encoder: early motion carries the hidden intent.",
        "Read the plot left to right: the encoder starts uncertain, then locks "
        "on as the prey commits. That matters because the planner can act "
        "before the final steps, while interception still has value. How it's "
        "generated: a supervised MLP classifier is trained on the prey's "
        "first-k absolute positions from oracle-condition rollouts, "
        "cross-validated; the bars are the capture ladder for the fixed "
        "predator conditions.")

    # 6 · hard intent ablation
    s = d.slide("This slide explains why the method carries a distribution. "
                "The wrong hard-intent condition is a stress test: 1.42 "
                "captures, far below 2.68 for the blind policy. Even the "
                "HONESTLY inferred hard intent at k=8 (2.56) loses to blind: "
                "the argmax throws away calibrated uncertainty. That failure "
                "is the argument for uncertainty-aware control.")
    d.title(s, "Hard intent labels can damage control", "Ablation")
    d.number(s, "2.68", "opponent-blind baseline", 0.9, 2.3, INK)
    d.number(s, "2.56", "hard intent inferred at k=8", 5.1, 2.3, INK)
    d.number(s, "1.42", "confident wrong intent", 9.3, 2.3, WARNING)
    d.text(s, "A point estimate throws away uncertainty. When the estimate is "
              "wrong, the predators commit to the wrong interception.",
           0.9, 5.2, 11.4, 1.0, size=20, bold=True, color=INK)

    # 7 · planner
    d.plot_slide(
        "The planner gets value from the whole belief", "Part 2", "part2_planner.png",
        ["Flat-belief planner: 3.07 captures/episode.",
         "Reactive oracle with true intent: 4.05.",
         "Planner with inferred belief: 4.31."],
        "The lift comes from inference plus lookahead, measured against a "
        "flat-belief ablation.",
        "Central evidence slide. The flat planner controls for generic "
        "lookahead, so the 3.07-to-4.31 step (+40%) is attributable to the "
        "opponent inference. Mechanics if asked: at each real step the "
        "planner enumerates all 125 joint predator actions, samples 8 intents "
        "from the belief (shared across candidates for a paired, low-variance "
        "comparison), rolls the true simulator 5 steps with the prey playing "
        "its policy under each sampled intent, and scores discounted captures "
        "minus distance to the imagined prey at the horizon: an interception.")

    # 8 · oracle wording
    s = d.slide("Defense slide. The planner has lookahead; the oracle policy "
                "reacts. That difference is why 4.31 can sit above 4.05 "
                "without overstating: say it exceeds a REACTIVE oracle, and "
                "name the oracle PLANNER as the true ceiling experiment still "
                "to run.")
    d.title(s, "The oracle comparison needs precise wording", "Interpretation")
    cols = [("Oracle policy", INK, "Receives the true intent, then reacts with "
             "the trained policy. Perfect intent information, no lookahead."),
            ("Belief planner", ACCENT, "Receives an inferred distribution, then "
             "searches short futures with simulator dynamics."),
            ("Reviewer line", WARNING, "Say it exceeds a reactive oracle. An "
             "oracle planner remains the real ceiling experiment.")]
    x = 0.7
    for t, c, body in cols:
        d.text(s, t, x, 1.8, 3.8, 0.5, size=19, bold=True, color=c)
        d.text(s, body, x, 2.35, 3.8, 2.4, size=15)
        x += 4.2
    d.text(s, "This wording preserves the strongest result without "
              "overstating the comparator.", 0.7, 5.4, 11.9, 0.6, size=18,
           bold=True, color=INK)

    # 9 · JEPA
    d.plot_slide(
        "JEPA recovers the belief without intent labels", "Representation",
        "jepa_vs_vae_encoder.png",
        ["Same trajectory data and the same 2-D latent size.",
         "JEPA probe 0.89 vs VAE probe 0.53 (3 encoder seeds).",
         "JEPA belief planner reaches 4.08 captures/episode."],
        "Predicting future behavior keeps the stable intent signal.",
        "Label-free route to the belief. JEPA predicts the representation of "
        "the future window through an EMA target encoder, no decoder; the "
        "prediction task rewards features that explain future motion, which "
        "IS the intent. The VAE spends capacity reconstructing evasion detail. "
        "Bar plot, if asked: probe = cross-validated logistic regression on "
        "the latent (is the strategy linearly decodable), ARI = unsupervised "
        "GMM clusters vs truth (does it cluster on its own); dashed line = "
        "supervised ceiling on the raw window, dotted = chance.")

    # 10 · learned dynamics negative
    d.plot_slide(
        "Learned dynamics failed the control test", "Boundary",
        "part4_jepa_world_planner.png",
        ["Simulator belief planning works.",
         "The latent world planner reaches 0.57 captures/episode.",
         "The control bottleneck is dynamics fidelity."],
        "Keep this negative result in the talk. It tells reviewers where the "
        "method breaks today.",
        "The honest boundary. An encoder can discard nuisance motion; a "
        "dynamics model must preserve the details that affect control, and "
        "the compressed latent loses them (5-step position error roughly "
        "3x the state-space model). Belief planning currently needs "
        "simulator dynamics.")

    # 11 · NEW: heatmaps ("circles are just squares")
    d.plot_slide(
        "“Are circles just squares?” Look at the occupancy", "Second axis",
        "occupancy_heatmaps.png",
        ["Corners specialist: four sharp hot-spots at the corners.",
         "Circle specialist: a diffuse, smeared ring, no crisp peaks.",
         "Different behaviours (cosine 0.42), but asymmetric: the circle "
         "signal is weak."],
        "The diffuse circle behaviour is why clustering fails and vanilla BC "
        "struggles; it motivates a genuinely distinct placement such as a "
        "snake path.",
        "Direct answer to the meeting question. Four resources on a circle "
        "are geometrically a small square of points, and the heatmap shows "
        "the consequence: the corners prey commits to four spots while the "
        "circle prey smears its time over a ring. Cosine similarity 0.42 and "
        "total-variation distance 0.47 between the two mean-occupancy "
        "distributions: different, but not cleanly separable. Generated from "
        "150 episodes per specialist, 100 steps each, 24x24 histogram of "
        "prey positions, resource geometry overlaid.")

    # 12 · NEW: scatter interpretation
    d.plot_slide(
        "The latent scatter: what the axes are", "Representation",
        "mopa_latent_resources_mappo_prey_occ.png",
        ["Each point = one episode's prey occupancy, embedded by the encoder; "
         "color = true placement (never shown to the encoder).",
         "Axes = the 2-D latent directly (PC1/PC2 of the latent when 4-D). "
         "No t-SNE, no rescaling.",
         "Linearly decodable (supervised 0.93, latents to 0.81) but ARI "
         "≈ 0: a direction, not two clusters."],
        "Signal yes, separated modes no. Shaping the latent (contrastive or "
        "equivariant) is the open problem.",
        "Generation pipeline, end to end: (1) roll out both specialists in "
        "their native envs, 1,200 episodes, 3 seeds, keep ABSOLUTE positions; "
        "(2) featurize each episode as a normalized 8x8 occupancy histogram, "
        "z-scored (a raw coordinate window is at chance even for a SUPERVISED "
        "probe, 0.54, so the featurization gates the task, not the encoder); "
        "(3) train VAE and JEPA unsupervised; (4) scatter the per-episode "
        "latent. The supervised probe is a labeled logistic regression on the "
        "same 64-D histogram: the ceiling. If asked which encoder wins here: "
        "at 2-D JEPA leads, at 4-D the VAE (0.81 vs 0.67), because occupancy "
        "already removed the noise JEPA's advantage feeds on.")

    # 13 · BC vs MAPPO
    d.plot_slide(
        "BC is strong enough to test strategy inputs", "Follow-up",
        "mopa_bc_vs_mappo.png",
        ["Random policy: 0.40 captures/episode.",
         "Vanilla BC: 1.22 captures/episode.",
         "MAPPO expert: 1.35 captures/episode."],
        "The clone recovers 86% of the expert edge, so the follow-up can test "
        "strategy information.",
        "Sanity check before the latent question: a poor clone would make the "
        "latent experiment unreadable. Deployed in the game against the same "
        "prey, the clone recovers 86% of the expert's edge over random, at an "
        "0.80 held-out action match (episode-level split). Cloning is not the "
        "bottleneck; the residual 14% is the headroom a strategy signal "
        "could fill.")

    # 14 · strategy info value (two panels)
    s = d.slide("Careful BC interpretation. Left, deployed captures: the "
                "placement-blind pooled clone gets 0.95; give it the TRUE "
                "placement and it reaches 1.11 vs expert 1.14, recovering 85% "
                "of the gap, so strategy information has real control value. "
                "The unsupervised latent recovers about 22%, within error "
                "bars today. Right, the observation sweep in action accuracy: "
                "as the encoder watches longer its probe climbs 0.60 to 0.78 "
                "and BC accuracy climbs with it, overtaking vanilla by 50 "
                "steps and reaching the oracle band by 125. Sell this as "
                "evidence a better strategy code matters, not as the main win.")
    d.title(s, "Strategy information has value, even when the latent is weak",
            "Follow-up")
    d.img(s, "mopa_bc_latent_deploy.png", 0.45, 1.5, 6.1, 4.4)
    d.img(s, "mopa_bc_latent_sweep.png", 6.8, 1.5, 6.1, 4.4)
    d.bullets(s, ["Oracle placement closes most of the deployed capture gap (85%).",
                  "The unsupervised latent gives a small deployed lift today (22%, within noise).",
                  "The sweep shows better probes translate into better BC."],
              0.7, 6.1, 12.0, 1.2, size=15)

    # 15 · NEW: averaging caveat
    s = d.slide("Get ahead of the sharpest question in the room. The "
                "placement is FIXED across a specialist's episodes, so when "
                "the encoder watches longer and the latent improves, that is "
                "a better ESTIMATE of a fixed quantity: statistical "
                "averaging, not a strategy unfolding. Contrast with the "
                "hidden-intent task, where the intent genuinely reveals "
                "itself within the episode. So describe the sweep as "
                "scaling-with-observation. What would prove more: a latent "
                "that clusters (contrastive or equivariant shaping), or a "
                "switching prey where the belief must track a change.")
    d.title(s, "Honest mechanism: is the gain just averaging?", "Caveat")
    d.bullets(s, [
        "The placement is fixed across a specialist's episodes.",
        "More observation → a better occupancy ESTIMATE: averaging, "
        "not the strategy unfolding.",
        "Different from the hidden-intent task, where intent is revealed "
        "within the episode.",
        "So: a scaling-with-observation result, stated without overclaim.",
        "Stronger tests next: a latent that clusters, or a switching prey."],
        1.0, 1.9, 11.3, 4.5, size=20)

    # 16 · NEW: where vanilla BC struggles
    s = d.slide("Concrete failure example; the GIF animates in slideshow "
                "mode. A held-out episode replays and each variant predicts "
                "every predator action: green ring = predicted the actual "
                "action, red = missed, with running accuracy. On "
                "route-dependent corners episodes the placement-blind clone "
                "cannot tell a committing prey from a passing one, hedges, "
                "and misses (0.60); told the placement, the same architecture "
                "follows the commit (0.83 circle, 0.92 corners). Caveat to "
                "volunteer: these are tail episodes chosen to show the "
                "effect; the average gap is the +2.4-point table.")
    d.title(s, "Where vanilla BC struggles", "Failure case")
    d.img(s, "demo_bc.gif", 2.4, 1.4, 8.5, 5.5)

    # 17 · NEW: how plots are generated
    s = d.slide("Rigor slide, keep it brief unless probed: every number is 3 "
                "seeds; all BC splits are episode-level so autocorrelated "
                "timesteps never straddle folds; conditioning latents only "
                "ever see steps strictly before the predicted episode; "
                "captures come from the raw +10-per-capture reward divided "
                "by 10, not the wrapper-scaled logs; and the repo has 12 "
                "fast tests plus a results checker that re-validates every "
                "paper-facing number against the raw npz files.")
    d.title(s, "How every plot is generated", "Procedures")
    d.bullets(s, [
        "3 seeds throughout; mean ± std shown on every bar and curve.",
        "Episode-level train/val splits; conditioning never sees the "
        "predicted steps (leakage-checked).",
        "Occupancy = normalized 8×8 histogram of prey positions; "
        "z-scored before any encoder.",
        "Supervised probe = labeled, cross-validated logistic regression on "
        "the same features: the ceiling.",
        "Captures = raw reward / 10 per first episode, deployed in the game.",
        "CI: 12 unit tests + mopa-check-results validates all paper-facing "
        "numbers against raw logs.",
        "Full procedures: docs/METHODS.md; interpretation guide on /present."],
        0.9, 1.8, 11.6, 4.8, size=17)

    # 18 · literature
    s = d.slide("Literature answer, three lanes. HOP (ICML 2024) couples a "
                "goal belief with tree search and is the closest "
                "system-level baseline to run first. The model-based MARL "
                "lane (MAZero, MAMBA, MARIE, MATWM) tests the planning and "
                "dynamics side. The opponent-modeling lane (MBOM, AORPO) "
                "tests whether an explicit belief beats conditioning on "
                "opponent history. For the BC thread the anchor is STRIL "
                "(AAAI 2025), per-trajectory strategy latents for imitation; "
                "for the clustering gap it is CTR (NeurIPS 2024), "
                "contrastive identity grounding.")
    d.title(s, "The literature slot is belief plus planning", "Positioning")
    cols = [("Closest comparator", ACCENT, "HOP couples goal belief with tree "
             "search. The most direct system-level baseline to run."),
            ("Model-based MARL", INK, "MAZero, MAMBA, MARIE, and MATWM test "
             "the planning and dynamics side of the story."),
            ("Opponent modeling", WARNING, "MBOM and AORPO test whether the "
             "opponent model adds value beyond policy conditioning.")]
    x = 0.7
    for t, c, body in cols:
        d.text(s, t, x, 1.8, 3.8, 0.5, size=19, bold=True, color=c)
        d.text(s, body, x, 2.35, 3.8, 2.2, size=15)
        x += 4.2
    d.text(s, "Paper line: calibrated opponent belief improves planning in a "
              "controlled hidden-intent task. Broader claims need adaptive "
              "opponents and external baselines.",
           0.7, 5.2, 12.0, 1.0, size=18, bold=True, color=INK)

    # 19-20 · anticipated questions (defend every plot)
    qa_pages = [
        [("Beating the oracle looks too good — is it real?",
          "It exceeds a REACTIVE oracle (true intent, no lookahead). Planning "
          "pre-positions, so 4.31 > 4.05. An oracle PLANNER is the true ceiling, "
          "and it's on the to-run list."),
         ("Why does an honestly inferred hard intent (2.56) lose to blind (2.68)?",
          "The argmax at k=8 is ~74% accurate; committing to a wrong corner "
          "costs more than hedging. That's the whole case for carrying the "
          "distribution, not a point estimate."),
         ("Probe vs ARI — what's the difference?",
          "Probe = can a LINEAR readout with labels decode the latent (chance "
          "0.25). ARI = does it cluster WITHOUT labels. On intent both are high; "
          "on circle-vs-corners only the probe is."),
         ("Isn't JEPA just a better-tuned VAE?",
          "Same data, latent size, and budget. JEPA predicts the future "
          "representation (no decoder); the VAE reconstructs. The gap (0.89 vs "
          "0.53) is the objective, not the tuning.")],
        [("Aren't circles just squares? (the meeting challenge)",
          "Geometrically yes — 4 points on a circle are a small square. "
          "Behaviourally no: corners commit to 4 sharp spots, the circle smears "
          "a diffuse ring (cosine 0.42). The circle signal is just weak."),
         ("If it's decodable (0.93), why doesn't it cluster (ARI≈0)?",
          "Linear separability ≠ two modes. The placement is a DIRECTION in "
          "latent space, not two blobs. Shaping it to cluster (contrastive / "
          "equivariant) is the open problem."),
         ("Is the latent-BC sweep gain just averaging over more episodes?",
          "Largely yes, and I say so: the placement is FIXED, so longer "
          "observation is a better occupancy ESTIMATE, not a strategy "
          "unfolding. It's a scaling-with-observation result."),
         ("The BC-failure GIF — cherry-picked?",
          "Deliberately: tail episodes chosen to show the mechanism. The "
          "average is the table (86% recovery, +2.4 accuracy pts). The GIF "
          "explains WHY, the table is the claim.")],
    ]
    for pi, page in enumerate(qa_pages):
        s = d.slide("Anticipated questions — answer crisply, then stop. These are "
                    "the sharpest ones the plots invite; each answer concedes "
                    "what's honest and states the claim precisely. Full set with "
                    "figures on the /present page's Figure guide.")
        d.title(s, "Questions they'll ask" + (" (cont.)" if pi else ""),
                "Defense" if pi == 0 else "Defense")
        y = 1.55
        for q, a in page:
            d.text(s, "Q  " + q, 0.7, y, 12.0, 0.5, size=16, bold=True, color=ACCENT)
            d.text(s, "A  " + a, 0.7, y + 0.5, 12.0, 0.8, size=14, color=MUTED)
            y += 1.42

    # 21 · close
    s = d.slide("Close on specific work. The repo has code hygiene, tests, "
                "result checks, and this deck. The scientific gap is the "
                "next experiment set: the oracle planner ceiling, a "
                "switching prey to test belief tracking, the HOP baseline, "
                "a genuinely distinct placement (snake path), and archived "
                "raw runs pinned to a release DOI. End with the concrete "
                "ask: accept the controlled result, fund the baseline and "
                "switching-opponent runs.")
    d.title(s, "Before submission, run the missing tests", "Close")
    rows = [("Oracle planner", "Measures the true planning ceiling with known intent."),
            ("Switching prey", "Tests whether the belief tracks a strategy change mid-episode."),
            ("HOP baseline", "Checks the closest belief-search comparator."),
            ("Snake placement", "A genuinely distinct behaviour class, so clusters can exist."),
            ("Raw artifact archive", "Pins checkpoints, logs, and metrics to a release DOI.")]
    y = 1.7
    for label, body in rows:
        d.text(s, label, 0.9, y, 3.1, 0.45, size=18, bold=True, color=INK)
        d.text(s, body, 4.2, y, 8.4, 0.55, size=17)
        y += 0.78
    d.text(s, "End with a concrete ask: accept the controlled result, then "
              "fund the baseline and switching-opponent run.",
           0.9, y + 0.25, 11.6, 0.6, size=18, bold=True, color=INK)

    out = os.path.join(ROOT, "marl_opp_aware_results.pptx")
    d.prs.save(out)
    print(f"saved {out} ({d.n} slides)")


if __name__ == "__main__":
    main()
